#!/usr/bin/python 
# -*- coding: utf-8 -*-
#author :hdlz
#version :v1.1
#describe:
'''
使用PPTX库生成产品信息Presentation的解析工具，
用以解决大量产品信息需要机械重复劳动的困境
1从xls数据中获取产品信息、图片地址链接、#价格策略
2生成具体的产品信息详情页
3具有简单的分类功能（等提出需求的时候再加入吧……）
'''
from pptx import Presentation
from pptx.util import Cm
from pptx.util import Pt
import os
import xlrd
from PIL import Image
im_v    =Image.open("竖版空白.jpg").resize((1920,1080))
im_h    =Image.open("横版空白.jpg")
#im_p    =Image.open("pic2.png")
a= 0
b=0
i=0
retext= '0'
#读取表格，加入选路径和选sheet索引的功能
xls_path=str(input("请放入xls文件地址：\n"))
xls_path = r'D:\lazyRocket\pptGO\rice_detail.xlsx'
workbook = xlrd.open_workbook(xls_path)
#输出索引号：
for i in workbook.sheet_names():
             print (workbook.sheet_names().index(i),i)
sheet_index = int(input("您想输出哪个sheet?\n"))
sheet = workbook.sheet_by_index(sheet_index)
print(os.getcwd())
filePath = str(os.getcwd())
filename = str(filePath+r'\路径.txt')


#这里使用了先前制作好的母版，1是class，2是detail#
SLD_class = 1
SLD_detail = 2

#导入母版
prs = Presentation("D:\lazyRocket\pptxGO\master_prs.pptx")

def riceDetail(a,b,c,p):
    '''
    idx 清单如下：
    0 Title 4
    13 Picture Placeholder 2
    15 Picture Placeholder 1
    17 Table Placeholder 3
    18 Table Placeholder 5
    a是标题，
    b是平面图链接
    c是效果图链接
    d是转置价测（未添加），
    e是rice_info表
    p是全表第几行的意思
    '''
    #生成详情页的专用版式
    title_slide_layout = prs.slide_layouts[SLD_detail]
    #写上标题
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = a
    

    try:
        print("正在绘制 ["+str(a)+"]  的图片部分")
        #插入图片，使用占位符的富文本功能，placeholde1    
        placeholder1 = slide.placeholders[13]
        placeholder1.name
        placeholder1.placeholder_format.type
        im_final = pic_resize(b).save('im_final.png')
        picture = placeholder1.insert_picture('im_final.png' )
        #插入图片2
        placeholder2 = slide.placeholders[15]
        placeholder2.name
        placeholder2.placeholder_format.type
        im_final = pic_resize(c).save('im_final.png')
        picture2 = placeholder2.insert_picture('im_final.png')
    except OSError:
        print("正在绘制 ["+str(a)+"]  的图片部分")
        #插入图片，使用占位符的富文本功能，placeholde1    
        placeholder1 = slide.placeholders[13]
        placeholder1.name
        placeholder1.placeholder_format.type
        im_final = pic_resize(b).save('im_final.jpg')
        picture = placeholder1.insert_picture('im_final.jpg' )
        #插入图片2
        placeholder2 = slide.placeholders[15]
        placeholder2.name
        placeholder2.placeholder_format.type
        im_final = pic_resize(c).save('im_final.jpg')
        picture2 = placeholder2.insert_picture('im_final.jpg')
    except Image.DecompressionBombError:
        with open(filename,'w') as f:
            f.write(str(a)+"\n")
    print("正在绘制  ["+str(a)+"]  的表格部分")
    #插入表格info
    
    table1 = slide.placeholders[18]
    #以这种方式插入的表具有原始占位符的位置和宽度。它的高度与行数成正比。
    #屈服了，我不搞了好吧，PPT似乎没法设置行高，我放弃使用模板中的格子了！
    
    table_info = table1.insert_table(rows= 11,cols = 2)
    table_info = slide.placeholders[18]
    table_info.table.cell(0,0).text = "产品信息"
    table_info.table.cell(0,1).text = "详情"
    
    for a in (0,1):
        for b in range(1,11):
            rice_table_info_index = [1,20,7,24,25,26,27,39,40,12,7,12,12,12]
            if a==0:
                table_info.table.cell(b,a).text = str(sheet.cell_value(0,rice_table_info_index[b]))
            else:
                table_info.table.cell(b,a).text = str(sheet.cell_value(p,rice_table_info_index[b]))
           # table_info.table.cell(b,a).font.size= Pt(14)
           # table_info.table.cell(b,a).font.name= "楷体"
    #插入表格价测
    
    table2 = slide.placeholders[17]
    #以这种方式插入的表具有原始占位符的位置和宽度。它的高度与行数成正比。
    #屈服了，我不搞了好吧，PPT似乎没法设置行高，我放弃使用模板中的格子了！

    table_price = table2.insert_table(rows= 6,cols = 2)
    table_price = slide.placeholders[17]
    table_price.table.cell(0,0).text = "项目"
    table_price.table.cell(0,1).text = "信息"
    
    for a in (0,1):
        for b in range(1,6):
            rice_table_price_index = [1,5,14,41,42,34,33,12,7,12,12,12]
            if a==0:
                table_price.table.cell(b,a).text = str(sheet.cell_value(0,rice_table_price_index[b]))
            else:
                table_price.table.cell(b,a).text = str(sheet.cell_value(p,rice_table_price_index[b]))
           # table_info.table.cell(b,a).font.size= Pt(14)

    #expection : Type error
def repeat():
#循环执行riceDetail()方法，传入坐标
    for i in range(1,sheet.nrows):
        
        riceDetail_title = sheet.cell_value(i,6)
        riceDetail_flat_url = sheet.cell_value(i,23)
        riceDetail_fact_url = sheet.cell_value(i,22)
        idx_sku = i

        riceDetail(riceDetail_title,riceDetail_flat_url ,riceDetail_fact_url,idx_sku)
#===========================
def pic_resize(a):
    im_p    =Image.open(a)
    if im_p.size[1] < im_p.size[0]:
        print("竖版图片")
        im_pic_resize =  resize_v(a) 
    else:
        print("横板图片")
        im_pic_resize = resize_v(a)#!!!!h变成了v
    return im_pic_resize
#竖版图片就得高度压缩到im_h的高度上
def resize_v(a):
    im_p = Image.open(a)
    cache = im_p.size[0]  * (im_v.size[1]/im_p.size[1])
    print(cache,im_v.size[1])
    im_f = im_p.resize((int(cache)  ,  int(im_v.size[1]) ))
#    im_f.show()
    return im_f
#横版图片就得宽度压缩到im_h的宽度上
def resize_h(a):
    im_p = Image.open(a)
    cache = im_p.size[1]  * (im_h.size[0]/im_p.size[0])
    print(im_h.size[0],cache)
    im_f = im_p.resize((int(im_v.size[0]),int(cache)))
#    im_f.show()
    return im_f

    mid_im_p_v = float(im_p.size[1])/2
#    print(mid_im_p_v )
    mid_im_p_h = float(im_p.size[0])/2
#    print(mid_im_p_h )

#========================================
#def kilogram_trans(text):
#    if '千克' in text:
#       retext =  text.replace('千克','KG')
#       return  retext
#    elif '0克' in text:
#       retext =  text.replace('0克','0G')
#       return  retext
    
def _main_():

#    im_final = pic_resize().save('im_final.jpg',img)
    repeat()

_main_()



#================
#保存最终成品  D:\lazyRocket\pptxGO\东北米1.xlsx
save_file_name = "信息看板"+input(
   '''请输入您想保存的名字:
    建议存储为当前时间用24h表示法的四位数字
    如1200、0859、1830等\n'''
    )+".pptx"
try:
    prs.save(save_file_name)
    print(save_file_name+" has been SAVED successful！\n")

except PermissionError:
    print("保存出错，请关闭"+str(save_file_name)+"重试\n")
